#!/usr/bin/env python3

"""
Author: Jasper Lelijveld
Date: 2019-01-11
Part of: AnimalFarm

Dependencies: https://python-pptx.readthedocs.io/en/latest/pip

Input: a filepath to a powerpoint presentation (file_location)
Output: a csv containing text from inside objects of the specified color (hex_color)

Args: filepath, hex color
"""

from pptx import Presentation
from csv import writer
import sys


def AntEater(file_location, hex_color):
    """
    :param file_location:
    :param hex_color:
    :return: csv
    """
    prs = Presentation(file_location)
    slideIndex = 0
    shapeText = []
    for slide in prs.slides:
        slideIndex += 1
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            fill = shape.fill
            fill.solid()
            try:
                if str(fill.fore_color.rgb) == str(hex_color):
                    print(shape.text)
                    shapeText.append([slideIndex, shape.text])
            except AttributeError:
                continue

    with open('output.csv', 'w', newline='') as csvfile:
        linewriter = writer(csvfile, delimiter=',')
        fieldnames = ['Slide number', 'Feedback']
        linewriter.writerow(fieldnames)
        for line in shapeText:
            linewriter.writerow(line)


if __name__ == '__main__':
    AntEater(*sys.argv[1:])
